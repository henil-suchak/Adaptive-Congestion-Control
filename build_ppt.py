"""Generator for the Adaptive-TCP-Congestion-Control academic presentation.

Run:  python build_ppt.py
Output: Adaptive_TCP_Congestion_Control.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree


# ── Theme ────────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0B, 0x2A, 0x4A)   # primary dark
BLUE       = RGBColor(0x1E, 0x4E, 0x8C)   # secondary
TEAL       = RGBColor(0x0E, 0x8A, 0xA6)   # accent
ORANGE     = RGBColor(0xE8, 0x7A, 0x17)   # highlight
GREEN      = RGBColor(0x1F, 0x8A, 0x4C)
RED        = RGBColor(0xC0, 0x39, 0x2B)
GRAY_DARK  = RGBColor(0x25, 0x2D, 0x3A)
GRAY_MED   = RGBColor(0x5E, 0x69, 0x7A)
GRAY_LIGHT = RGBColor(0xE6, 0xEA, 0xF0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG         = RGBColor(0xF7, 0xF9, 0xFC)   # page background
CODE_BG    = RGBColor(0x10, 0x17, 0x25)
CODE_FG    = RGBColor(0xE6, 0xEA, 0xF0)

FONT_TITLE = "Calibri"
FONT_BODY  = "Calibri"
FONT_MONO  = "Consolas"


# ── Presentation setup (16:9) ────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    if not shadow:
        shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.12
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=GRAY_DARK,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_multirun(slide, x, y, w, h, runs, *, size=14, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """runs = list of (text, {bold, color, italic, font, size})."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        r.font.name = opts.get("font", FONT_BODY)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", GRAY_DARK)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=GRAY_DARK,
                bullet_color=TEAL, line_spacing=1.35, bold_first=False,
                bullet_char="●"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet glyph
        r_b = p.add_run()
        r_b.text = f"{bullet_char}  "
        r_b.font.name = FONT_BODY
        r_b.font.size = Pt(size)
        r_b.font.bold = True
        r_b.font.color.rgb = bullet_color
        # content (may contain a bold-colon split)
        if isinstance(item, tuple) and len(item) == 2:
            head, tail = item
            r1 = p.add_run()
            r1.text = head
            r1.font.name = FONT_BODY
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = tail
            r2.font.name = FONT_BODY
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        elif isinstance(item, tuple) and len(item) == 1:
            r = p.add_run()
            r.text = item[0]
            r.font.name = FONT_BODY
            r.font.size = Pt(size)
            r.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = FONT_BODY
            r.font.size = Pt(size)
            r.font.bold = bold_first and i == 0
            r.font.color.rgb = color
    return tb


def slide_background(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # push to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def slide_header(slide, title, kicker=None, page_num=None, total=None):
    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background(); bar.shadow.inherit = False
    # Title
    add_text(slide, Inches(0.55), Inches(0.32), Inches(10.5), Inches(0.7),
             title, size=28, bold=True, color=NAVY)
    if kicker:
        add_text(slide, Inches(0.55), Inches(0.92), Inches(10.5), Inches(0.35),
                 kicker, size=14, color=TEAL, italic=True)
    # Divider
    div = slide.shapes.add_connector(1, Inches(0.55), Inches(1.35),
                                     Inches(12.8), Inches(1.35))
    div.line.color.rgb = GRAY_LIGHT
    div.line.width = Pt(1.25)
    # Footer
    add_text(slide, Inches(0.55), Inches(7.05), Inches(8),
             Inches(0.3), "Adaptive TCP Congestion Control using Deep RL",
             size=9, color=GRAY_MED)
    if page_num is not None:
        add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5),
                 Inches(0.3), f"{page_num} / {total}",
                 size=9, color=GRAY_MED, align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, header, rows, *,
              header_fill=NAVY, header_fg=WHITE,
              stripe=GRAY_LIGHT, body_fg=GRAY_DARK,
              first_col_bold=False, col_colors=None, header_size=12,
              body_size=11):
    shape = slide.shapes.add_table(len(rows) + 1, len(header), x, y, w, h)
    tbl = shape.table

    # header
    for j, cell_text in enumerate(header):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.margin_left = cell.margin_right = Inches(0.1)
        cell.margin_top = cell.margin_bottom = Inches(0.04)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = cell_text
        r.font.name = FONT_BODY
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = header_fg
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # body
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (WHITE if i % 2 == 0 else stripe)
            # optional per-cell override
            fg = body_fg
            bold = first_col_bold and j == 0
            if col_colors and j in col_colors and col_colors[j].get(i):
                fg = col_colors[j][i]
                bold = True
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.name = FONT_BODY
            r.font.size = Pt(body_size)
            r.font.bold = bold
            r.font.color.rgb = fg
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl


def add_code_block(slide, x, y, w, h, code, title=None):
    if title:
        add_text(slide, x, y, w, Inches(0.3), title, size=11,
                 bold=True, color=TEAL)
        y = y + Inches(0.32)
        h = h - Inches(0.32)
    bg = add_rounded(slide, x, y, w, h, CODE_BG)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        r.font.name = FONT_MONO
        r.font.size = Pt(11)
        r.font.color.rgb = CODE_FG


def add_arrow(slide, x1, y1, x2, y2, color=TEAL, weight=2.0):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    # arrow end
    ln = conn.line._get_or_add_ln()
    tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")
    return conn


def add_pill(slide, x, y, w, h, text, *, fill=TEAL, fg=WHITE, size=11,
             bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.5
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background(); shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = fg
    return shp


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE CONSTRUCTION
# ─────────────────────────────────────────────────────────────────────────────

TOTAL_SLIDES = 27  # final count


# ── Slide 1 — Title ──────────────────────────────────────────────────────────
def build_title():
    s = prs.slides.add_slide(BLANK)
    # gradient-like bands
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, 0, SW, Inches(3.8), BLUE)
    # left accent strip
    add_rect(s, 0, Inches(3.8), SW, Inches(0.08), TEAL)

    # Small badge
    add_rounded(s, Inches(0.7), Inches(0.7), Inches(3.3), Inches(0.45), TEAL)
    add_text(s, Inches(0.7), Inches(0.7), Inches(3.3), Inches(0.45),
             "B.TECH CE  •  SEMESTER VI  •  PROJECT",
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # Title
    add_text(s, Inches(0.7), Inches(1.6), Inches(12), Inches(1.3),
             "Adaptive TCP Congestion Control",
             size=46, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(2.5), Inches(12), Inches(1.0),
             "using Deep Reinforcement Learning",
             size=30, color=RGBColor(0xB8, 0xD4, 0xF5))

    # Divider tick
    add_rect(s, Inches(0.7), Inches(3.55), Inches(1.2), Inches(0.05), TEAL)

    # Team
    add_text(s, Inches(0.7), Inches(4.1), Inches(5), Inches(0.35),
             "PRESENTED BY",
             size=11, bold=True, color=TEAL)
    team = [
        "Jewel Bhannvadiya",
        "Archi Patel",
        "Henil Suchak",
    ]
    ids = ["CE-065", "CE-087", "CE-148"]
    for i, (name, code) in enumerate(zip(team, ids)):
        y = Inches(4.5 + i * 0.42)
        add_text(s, Inches(0.7), y, Inches(4), Inches(0.4),
                 name, size=16, bold=True, color=NAVY)
        add_pill(s, Inches(3.35), y + Inches(0.06), Inches(1.0),
                 Inches(0.3), code, fill=GRAY_LIGHT, fg=NAVY, size=10)

    # Right-hand card (guide / subject)
    card_x = Inches(7.7); card_y = Inches(4.1)
    card_w = Inches(5.0); card_h = Inches(2.7)
    add_rounded(s, card_x, card_y, card_w, card_h, WHITE)
    add_rect(s, card_x, card_y, Inches(0.08), card_h, ORANGE)
    add_text(s, card_x + Inches(0.3), card_y + Inches(0.2),
             card_w - Inches(0.5), Inches(0.35),
             "GUIDED BY", size=11, bold=True, color=TEAL)
    add_text(s, card_x + Inches(0.3), card_y + Inches(0.55),
             card_w - Inches(0.5), Inches(0.5),
             "Prof. Bhavika M. Gambhava",
             size=18, bold=True, color=NAVY)
    add_text(s, card_x + Inches(0.3), card_y + Inches(1.1),
             card_w - Inches(0.5), Inches(0.35),
             "SUBJECT", size=11, bold=True, color=TEAL)
    add_text(s, card_x + Inches(0.3), card_y + Inches(1.45),
             card_w - Inches(0.5), Inches(0.45),
             "System Design Practice  (23CE625)",
             size=14, bold=True, color=GRAY_DARK)
    add_text(s, card_x + Inches(0.3), card_y + Inches(2.0),
             card_w - Inches(0.5), Inches(0.35),
             "Department of Computer Engineering",
             size=11, color=GRAY_DARK)
    add_text(s, card_x + Inches(0.3), card_y + Inches(2.28),
             card_w - Inches(0.5), Inches(0.35),
             "Dharmsinh Desai University",
             size=11, italic=True, color=GRAY_MED)


# ── Slide 2 — Agenda ─────────────────────────────────────────────────────────
def build_agenda():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Agenda", kicker="What we'll walk through today",
                 page_num=2, total=TOTAL_SLIDES)

    items = [
        ("Problem & Motivation", "Why congestion control matters today"),
        ("TCP Fundamentals",     "cwnd, RTT, packet loss in plain terms"),
        ("Limitations of Classical Algorithms", "CUBIC, Reno and their blind spots"),
        ("Reinforcement Learning Approach", "Mapping TCP to an MDP"),
        ("Soft Actor-Critic (SAC)", "Why this algorithm and how it learns"),
        ("System Architecture", "ns-3 • Python • Spring Boot • React"),
        ("Implementation Details", "Shared memory, training, reward design"),
        ("Live Dashboard Demo", "Real-time monitoring of SAC vs CUBIC"),
        ("Conclusion & Future Work", "What we achieved and what's next"),
    ]
    # Two columns of cards
    left_x = Inches(0.7); right_x = Inches(6.9)
    card_w = Inches(5.7); card_h = Inches(0.9)
    for i, (title, sub) in enumerate(items):
        col = i % 2
        row = i // 2
        x = left_x if col == 0 else right_x
        y = Inches(1.6 + row * 1.05)
        add_rounded(s, x, y, card_w, card_h, WHITE)
        # number badge
        add_rounded(s, x + Inches(0.15), y + Inches(0.15),
                    Inches(0.6), Inches(0.6), NAVY)
        add_text(s, x + Inches(0.15), y + Inches(0.15),
                 Inches(0.6), Inches(0.6), f"{i+1:02d}",
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.95), y + Inches(0.1),
                 card_w - Inches(1.1), Inches(0.4),
                 title, size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.95), y + Inches(0.48),
                 card_w - Inches(1.1), Inches(0.4),
                 sub, size=11, color=GRAY_MED)


# ── Slide 3 — Problem Statement ──────────────────────────────────────────────
def build_problem():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Problem Statement", kicker="The core challenge in TCP",
                 page_num=3, total=TOTAL_SLIDES)

    # Quote / big question card
    add_rounded(s, Inches(0.7), Inches(1.65), Inches(12), Inches(1.2), NAVY)
    add_text(s, Inches(1.0), Inches(1.78), Inches(11.4), Inches(0.45),
             "THE CORE QUESTION", size=11, bold=True, color=TEAL)
    add_text(s, Inches(1.0), Inches(2.10), Inches(11.4), Inches(0.7),
             "“How fast should a sender transmit data without overwhelming the network?”",
             size=20, bold=True, color=WHITE, italic=True)

    # Limitations column
    add_text(s, Inches(0.7), Inches(3.1), Inches(5.8), Inches(0.4),
             "LIMITATIONS OF CLASSICAL ALGORITHMS",
             size=12, bold=True, color=ORANGE)
    add_bullets(s, Inches(0.7), Inches(3.55), Inches(5.8), Inches(3.3), [
        ("Fixed rules: ",  "CUBIC / Reno use hand-tuned formulas"),
        ("React, don't predict: ",  "only respond after packet loss"),
        ("Single signal: ",  "loss is the only input used"),
        ("Static: ",  "one algorithm fits every network type"),
        ("Loss-blind: ",  "can't tell random loss from congestion"),
    ], size=13)

    # Our goal column
    add_text(s, Inches(6.9), Inches(3.1), Inches(5.8), Inches(0.4),
             "OUR GOAL",
             size=12, bold=True, color=TEAL)
    add_rounded(s, Inches(6.9), Inches(3.55), Inches(5.8), Inches(3.3), WHITE)
    add_rect(s, Inches(6.9), Inches(3.55), Inches(5.8), Inches(0.07), TEAL)
    add_text(s, Inches(7.15), Inches(3.75), Inches(5.3), Inches(0.5),
             "Learning-based adaptive controller",
             size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(7.15), Inches(4.25), Inches(5.3), Inches(2.5), [
        "Learns from experience, not hand-coded rules",
        "Uses multiple signals: RTT, throughput, cwnd, loss",
        "Adapts in real-time to changing conditions",
        "Robust against non-congestion packet loss",
    ], size=12)


# ── Slide 4 — Why This Matters ───────────────────────────────────────────────
def build_why_matters():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Why Does This Matter?",
                 kicker="Poor congestion control hurts real users every day",
                 page_num=4, total=TOTAL_SLIDES)

    # Left: app impact tiles
    apps = [
        ("📺", "Video Streaming",  "Buffering, quality drops"),
        ("📞", "Video Calls",      "Latency spikes, stutter"),
        ("🎮", "Online Gaming",    "Lag, disconnections"),
        ("⬇", "Large Downloads",  "Wasted bandwidth"),
        ("🏢", "Data Centers",     "Incast collapse"),
        ("📡", "5G / Satellite",   "Highly dynamic links"),
    ]
    grid_x = Inches(0.7); grid_y = Inches(1.65)
    tile_w = Inches(3.9); tile_h = Inches(1.2)
    gap = Inches(0.25)
    for i, (icon, name, desc) in enumerate(apps):
        col = i % 3; row = i // 3
        x = grid_x + col * (tile_w + gap)
        y = grid_y + row * (tile_h + gap)
        add_rounded(s, x, y, tile_w, tile_h, WHITE)
        add_rect(s, x, y, Inches(0.08), tile_h, TEAL)
        add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(0.7),
                 Inches(0.7), icon, size=26, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.95), y + Inches(0.15),
                 tile_w - Inches(1.0), Inches(0.45),
                 name, size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.95), y + Inches(0.6),
                 tile_w - Inches(1.0), Inches(0.5),
                 desc, size=11, color=GRAY_MED)

    # Insight banner
    add_rounded(s, Inches(0.7), Inches(5.6), Inches(12.0), Inches(1.2), NAVY)
    add_text(s, Inches(1.0), Inches(5.75), Inches(11.5), Inches(0.4),
             "KEY INSIGHT",
             size=11, bold=True, color=ORANGE)
    add_text(s, Inches(1.0), Inches(6.05), Inches(11.5), Inches(0.7),
             "Networks today are highly diverse — 5G, satellite, Wi-Fi, data-center fabrics.",
             size=15, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(6.40), Inches(11.5), Inches(0.5),
             "A single fixed rule cannot perform optimally across all of them.",
             size=13, color=RGBColor(0xB8, 0xD4, 0xF5), italic=True)


# ── Slide 5 — What is TCP? ───────────────────────────────────────────────────
def build_what_is_tcp():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Background — What is TCP?",
                 kicker="The protocol that carries most of the internet",
                 page_num=5, total=TOTAL_SLIDES)

    # TCP services row (4 tiles)
    tiles = [
        ("Reliability",     "Lost packets are retransmitted"),
        ("Ordering",        "Packets arrive in sequence"),
        ("Flow Control",    "Don't overwhelm the receiver"),
        ("Congestion Ctrl", "Don't overwhelm the network"),
    ]
    x0 = Inches(0.7); y0 = Inches(1.65)
    tw = Inches(2.95); th = Inches(1.3); gap = Inches(0.15)
    for i, (t, d) in enumerate(tiles):
        x = x0 + i * (tw + gap)
        add_rounded(s, x, y0, tw, th, WHITE)
        add_rect(s, x, y0, tw, Inches(0.08), TEAL if i < 3 else ORANGE)
        add_text(s, x + Inches(0.2), y0 + Inches(0.2),
                 tw - Inches(0.4), Inches(0.45),
                 t, size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.2), y0 + Inches(0.68),
                 tw - Inches(0.4), Inches(0.5),
                 d, size=11, color=GRAY_MED)

    # Congestion control highlighted
    add_text(s, Inches(0.7), Inches(3.2), Inches(12), Inches(0.4),
             "THE KEY VARIABLE",
             size=12, bold=True, color=ORANGE)

    # Equation card
    add_rounded(s, Inches(0.7), Inches(3.6), Inches(12.0), Inches(1.2), NAVY)
    add_text(s, Inches(1.0), Inches(3.75), Inches(11.5), Inches(0.4),
             "Congestion Window  (cwnd)",
             size=13, bold=True, color=TEAL)
    add_text(s, Inches(1.0), Inches(4.05), Inches(11.5), Inches(0.7),
             "cwnd  =  bytes allowed “in flight” (sent but not yet acknowledged)",
             size=18, bold=True, color=WHITE, font=FONT_MONO)

    # Explanation
    add_bullets(s, Inches(0.7), Inches(5.05), Inches(12), Inches(1.7), [
        ("cwnd grows ", "on every successful ACK — sender can push faster"),
        ("cwnd shrinks ", "on packet loss — sender must slow down"),
        ("The entire congestion-control game is: ",
         "how to grow cwnd and how much to cut it on loss"),
    ], size=13)


# ── Slide 6 — How Congestion Happens ─────────────────────────────────────────
def build_congestion():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "How Congestion Happens",
                 kicker="Bottleneck links, finite queues, and packet loss",
                 page_num=6, total=TOTAL_SLIDES)

    # Diagram — sender, router+buffer, receiver
    y_mid = Inches(2.8)
    # sender
    add_rounded(s, Inches(0.8), y_mid, Inches(2.2), Inches(1.2), TEAL)
    add_text(s, Inches(0.8), y_mid, Inches(2.2), Inches(1.2),
             "Sender\n(fast)",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # router
    add_rounded(s, Inches(4.5), y_mid, Inches(3.5), Inches(1.2), NAVY)
    add_text(s, Inches(4.5), y_mid, Inches(3.5), Inches(1.2),
             "Router\n(bottleneck link)",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # queue under router
    add_rounded(s, Inches(4.5), y_mid + Inches(1.4),
                Inches(3.5), Inches(0.8), ORANGE)
    add_text(s, Inches(4.5), y_mid + Inches(1.4),
             Inches(3.5), Inches(0.8),
             "Queue fills up  →  packets drop",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # receiver
    add_rounded(s, Inches(9.5), y_mid, Inches(2.2), Inches(1.2), GREEN)
    add_text(s, Inches(9.5), y_mid, Inches(2.2), Inches(1.2),
             "Receiver\n(slow link)",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # arrows
    add_arrow(s, Inches(3.05), y_mid + Inches(0.6),
              Inches(4.45), y_mid + Inches(0.6), color=GRAY_DARK, weight=2.5)
    add_arrow(s, Inches(8.05), y_mid + Inches(0.6),
              Inches(9.45), y_mid + Inches(0.6), color=GRAY_DARK, weight=2.5)

    # Symptoms list below
    add_text(s, Inches(0.7), Inches(5.2), Inches(12), Inches(0.4),
             "WHAT THE SENDER OBSERVES",
             size=12, bold=True, color=ORANGE)
    sx = Inches(0.7); sy = Inches(5.6)
    tw = Inches(3.95); th = Inches(1.2); gap = Inches(0.15)
    syms = [
        ("RTT rises", "Packets queued → longer round-trip time"),
        ("Loss appears", "Queue overflows → router drops packets"),
        ("Throughput stalls", "Link saturated → no gain from sending faster"),
    ]
    for i, (t, d) in enumerate(syms):
        x = sx + i * (tw + gap)
        add_rounded(s, x, sy, tw, th, WHITE)
        add_rect(s, x, sy, Inches(0.08), th, RED)
        add_text(s, x + Inches(0.2), sy + Inches(0.15),
                 tw - Inches(0.35), Inches(0.4), t,
                 size=13, bold=True, color=RED)
        add_text(s, x + Inches(0.2), sy + Inches(0.55),
                 tw - Inches(0.35), Inches(0.6), d,
                 size=11, color=GRAY_DARK)


# ── Slide 7 — Classical Algorithms ───────────────────────────────────────────
def build_classical():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Traditional Algorithms — Reno & CUBIC",
                 kicker="Hand-tuned rules that have served us since 1990",
                 page_num=7, total=TOTAL_SLIDES)

    # Two cards side by side
    card_w = Inches(5.85); card_h = Inches(4.0)
    # Reno
    rx, ry = Inches(0.7), Inches(1.65)
    add_rounded(s, rx, ry, card_w, card_h, WHITE)
    add_rect(s, rx, ry, card_w, Inches(0.55), BLUE)
    add_text(s, rx + Inches(0.3), ry, card_w, Inches(0.55),
             "TCP Reno  (1990)", size=16, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, rx + Inches(0.3), ry + Inches(0.75),
             card_w - Inches(0.6), Inches(0.4),
             "Additive Increase / Multiplicative Decrease",
             size=12, bold=True, color=TEAL, italic=True)
    add_bullets(s, rx + Inches(0.3), ry + Inches(1.2),
                card_w - Inches(0.6), Inches(2.5), [
        ("Increase: ", "cwnd += 1 segment per RTT"),
        ("On loss:  ", "cwnd  =  cwnd / 2"),
        ("Simple, predictable sawtooth pattern",),
        ("Too slow on high-bandwidth × delay links",),
    ], size=12)

    # CUBIC
    cx = Inches(6.8)
    add_rounded(s, cx, ry, card_w, card_h, WHITE)
    add_rect(s, cx, ry, card_w, Inches(0.55), NAVY)
    add_text(s, cx + Inches(0.3), ry, card_w, Inches(0.55),
             "TCP CUBIC  (Linux default, 2006)",
             size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx + Inches(0.3), ry + Inches(0.75),
             card_w - Inches(0.6), Inches(0.4),
             "cwnd grows as a cubic function of time since last loss",
             size=12, bold=True, color=TEAL, italic=True)
    # Formula block
    add_code_block(s, cx + Inches(0.3), ry + Inches(1.2),
                   card_w - Inches(0.6), Inches(0.7),
                   "W(t)  =  C · (t − K)³  +  W_max")
    add_bullets(s, cx + Inches(0.3), ry + Inches(2.05),
                card_w - Inches(0.6), Inches(1.85), [
        ("RTT-independent growth — fair across delays",),
        ("Fast recovery after loss",),
        ("Still reacts only after loss occurs",),
    ], size=12)

    # Key problem banner
    add_rounded(s, Inches(0.7), Inches(5.85), Inches(12.0), Inches(0.95), NAVY)
    add_text(s, Inches(1.0), Inches(5.95), Inches(11.5), Inches(0.4),
             "COMMON LIMITATION",
             size=11, bold=True, color=ORANGE)
    add_text(s, Inches(1.0), Inches(6.25), Inches(11.5), Inches(0.5),
             "Both react only AFTER loss — by then the queue is already full.",
             size=14, bold=True, color=WHITE)


# ── Slide 8 — CUBIC Sawtooth ─────────────────────────────────────────────────
def build_sawtooth():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "The CUBIC Sawtooth Pattern",
                 kicker="How cwnd behaves over time — and what's wrong with it",
                 page_num=8, total=TOTAL_SLIDES)

    # Draw sawtooth chart area
    cx, cy, cw, ch = Inches(0.7), Inches(1.7), Inches(7.5), Inches(4.8)
    add_rounded(s, cx, cy, cw, ch, WHITE)
    # Axes
    axis_l = cx + Inches(0.7); axis_b = cy + ch - Inches(0.6)
    axis_t = cy + Inches(0.5); axis_r = cx + cw - Inches(0.3)
    # y-axis
    add_rect(s, axis_l, axis_t, Emu(6000), axis_b - axis_t, GRAY_DARK)
    # x-axis
    add_rect(s, axis_l, axis_b, axis_r - axis_l, Emu(6000), GRAY_DARK)
    # labels
    add_text(s, cx + Inches(0.1), cy + Inches(0.15), Inches(1.0), Inches(0.3),
             "cwnd", size=11, bold=True, color=NAVY)
    add_text(s, axis_r - Inches(0.4), axis_b + Inches(0.08),
             Inches(0.7), Inches(0.3), "time",
             size=11, bold=True, color=NAVY)

    # Draw teeth as three triangles
    import math
    tooth_w = (axis_r - axis_l) / 3
    for i in range(3):
        start_x = axis_l + i * tooth_w
        # growing cubic-ish (rising polyline approximated with 2 lines)
        peak_x = start_x + tooth_w * 0.85
        peak_y = axis_t + Inches(0.3)
        mid_x  = start_x + tooth_w * 0.4
        mid_y  = axis_t + Inches(2.0)
        # rising slow → fast (two segments: concave, convex)
        l1 = s.shapes.add_connector(1, start_x, axis_b, mid_x, mid_y)
        l1.line.color.rgb = BLUE; l1.line.width = Pt(3.0)
        l2 = s.shapes.add_connector(1, mid_x, mid_y, peak_x, peak_y)
        l2.line.color.rgb = BLUE; l2.line.width = Pt(3.0)
        # vertical drop (loss)
        drop_end_y = axis_t + Inches(2.5)
        l3 = s.shapes.add_connector(1, peak_x, peak_y, peak_x, drop_end_y)
        l3.line.color.rgb = RED; l3.line.width = Pt(3.0)
        # loss marker
        lm = s.shapes.add_shape(MSO_SHAPE.OVAL, peak_x - Inches(0.07),
                                peak_y - Inches(0.07), Inches(0.14),
                                Inches(0.14))
        lm.fill.solid(); lm.fill.fore_color.rgb = RED
        lm.line.fill.background(); lm.shadow.inherit = False
        # connect from drop to start of next tooth
        next_x = start_x + tooth_w
        l4 = s.shapes.add_connector(1, peak_x, drop_end_y, next_x, axis_b)
        l4.line.color.rgb = BLUE; l4.line.width = Pt(3.0)
        # loss label
        add_text(s, peak_x - Inches(0.35), peak_y - Inches(0.5),
                 Inches(0.9), Inches(0.3), "LOSS",
                 size=9, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # Right column — characteristics
    lx = Inches(8.5); ly = Inches(1.7)
    add_text(s, lx, ly, Inches(4.5), Inches(0.4),
             "CHARACTERISTICS", size=12, bold=True, color=ORANGE)
    add_bullets(s, lx, ly + Inches(0.45), Inches(4.5), Inches(4.5), [
        ("Aggressive growth: ", "cubic curve probes bandwidth fast"),
        ("Sharp drop: ",         "cwnd ×0.7 on every loss event"),
        ("Loss-triggered: ",     "only reacts after damage is done"),
        ("Loss-blind: ",         "can't tell random from congestion loss"),
        ("Bufferbloat: ",        "RTT rises steadily before each drop"),
    ], size=12)


# ── Slide 9 — Why RL ─────────────────────────────────────────────────────────
def build_why_rl():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Why Reinforcement Learning?",
                 kicker="TCP congestion control is a textbook sequential decision problem",
                 page_num=9, total=TOTAL_SLIDES)

    # Mapping table
    rows = [
        ["Agent",       "Decision maker",                  "The TCP sender"],
        ["Environment", "World the agent acts in",         "The network (ns-3)"],
        ["State",       "What the agent observes",         "RTT, throughput, cwnd, loss"],
        ["Action",      "What the agent does",             "Adjust cwnd (scale factor)"],
        ["Reward",      "Scalar feedback signal",          "High throughput, low RTT, no loss"],
    ]
    add_table(s, Inches(0.7), Inches(1.65), Inches(7.6), Inches(3.3),
              ["RL Concept", "Meaning", "TCP Mapping"], rows,
              first_col_bold=True, header_size=12, body_size=12)

    # Advantages column (right)
    ax = Inches(8.6); ay = Inches(1.65)
    add_text(s, ax, ay, Inches(4.2), Inches(0.4),
             "KEY ADVANTAGES", size=12, bold=True, color=ORANGE)
    add_bullets(s, ax, ay + Inches(0.45), Inches(4.2), Inches(4.5), [
        "Learns optimal behavior from experience",
        "Combines multiple signals, not just loss",
        "Can anticipate congestion — act before loss",
        "Generalizes to unseen network conditions",
        "Improves automatically with more data",
    ], size=12)

    # Bottom banner
    add_rounded(s, Inches(0.7), Inches(5.6), Inches(12.0), Inches(1.2), NAVY)
    add_text(s, Inches(1.0), Inches(5.75), Inches(11.5), Inches(0.4),
             "THE INTUITION",
             size=11, bold=True, color=TEAL)
    add_text(s, Inches(1.0), Inches(6.05), Inches(11.5), Inches(0.5),
             "Instead of hand-coding rules, the agent discovers them by trial and feedback.",
             size=15, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(6.38), Inches(11.5), Inches(0.4),
             "Good actions raise reward; bad ones lower it. The policy drifts toward good.",
             size=12, color=RGBColor(0xB8, 0xD4, 0xF5), italic=True)


# ── Slide 10 — MDP Formulation ───────────────────────────────────────────────
def build_mdp():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Our MDP Formulation",
                 kicker="The exact state, action, and normalization we use",
                 page_num=10, total=TOTAL_SLIDES)

    # State space table
    add_text(s, Inches(0.7), Inches(1.6), Inches(7.6), Inches(0.4),
             "STATE SPACE  —  6 features, normalized to [0, 1]",
             size=12, bold=True, color=TEAL)
    rows = [
        ["cwnd",          "Congestion window (bytes)",    "1,400,000"],
        ["rtt_us",        "Round-trip time (µs)",         "200,000"],
        ["throughput",    "Bytes / second",               "250,000"],
        ["packetLoss",    "Loss count in interval",       "100"],
        ["segmentSize",   "TCP segment size (bytes)",     "1,500"],
        ["bytesInFlight", "Outstanding unACKed bytes",    "1,400,000"],
    ]
    add_table(s, Inches(0.7), Inches(2.05), Inches(7.6), Inches(3.4),
              ["Feature", "Description", "Max Value"], rows,
              first_col_bold=True, header_size=12, body_size=11)

    # Right — action space card
    ax = Inches(8.55); ay = Inches(1.6)
    add_text(s, ax, ay, Inches(4.3), Inches(0.4),
             "ACTION SPACE", size=12, bold=True, color=ORANGE)
    add_rounded(s, ax, ay + Inches(0.45), Inches(4.3), Inches(1.3), WHITE)
    add_rect(s, ax, ay + Inches(0.45), Inches(0.08), Inches(1.3), ORANGE)
    add_text(s, ax + Inches(0.2), ay + Inches(0.55),
             Inches(4.0), Inches(0.45),
             "Continuous factor ∈ [ 0.8, 1.2 ]",
             size=14, bold=True, color=NAVY)
    add_text(s, ax + Inches(0.2), ay + Inches(1.0),
             Inches(4.0), Inches(0.75),
             "Scales cwnd by ±20 % each step —\nsmooth, scale-free control.",
             size=11, color=GRAY_DARK, italic=True)

    # Formula card
    add_rounded(s, ax, ay + Inches(1.95), Inches(4.3), Inches(1.2), NAVY)
    add_text(s, ax + Inches(0.2), ay + Inches(2.05),
             Inches(4.0), Inches(0.3),
             "UPDATE RULE", size=10, bold=True, color=TEAL)
    add_text(s, ax + Inches(0.2), ay + Inches(2.35),
             Inches(4.0), Inches(0.8),
             "new_cwnd = cwnd × action",
             size=14, bold=True, color=WHITE, font=FONT_MONO)

    # Step interval pill
    add_rounded(s, ax, ay + Inches(3.35), Inches(4.3), Inches(0.6), BLUE)
    add_text(s, ax, ay + Inches(3.35), Inches(4.3), Inches(0.6),
             "Step interval:  10 ms (simulated time)",
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Footer note
    add_text(s, Inches(0.7), Inches(5.7), Inches(12.0), Inches(1.2),
             "Normalization ensures every feature contributes on the same scale — "
             "without it, the larger-magnitude inputs (e.g. cwnd in bytes) would dominate "
             "the neural network and small-but-important features (e.g. loss count) "
             "would be ignored.",
             size=12, italic=True, color=GRAY_MED)


# ── Slide 11 — Reward Function ───────────────────────────────────────────────
def build_reward():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Reward Function Design",
                 kicker="A single scalar that balances throughput, latency and loss",
                 page_num=11, total=TOTAL_SLIDES)

    # Components as 4 cards with signs
    parts = [
        ("+", "Throughput term",  "√(throughput / T_max)",
         "Concave — diminishing returns stop the agent from being greedy",
         GREEN),
        ("−", "RTT penalty",      "((RTT − RTT_min) / 400k)² × 40",
         "Quadratic — strongly discourages buffer bloat",
         RED),
        ("−", "Loss penalty",     "0.3 + 0.05 × loss",
         "Cliff penalty — discrete, triggered on any loss event",
         RED),
        ("+", "Stability bonus",  "0.1 × Gaussian(cwnd ≈ BDP)",
         "Nudges the policy toward the ideal operating point",
         BLUE),
    ]
    x0 = Inches(0.7); y0 = Inches(1.65); cw = Inches(6.0); ch = Inches(1.1)
    for i, (sign, title, formula, desc, col) in enumerate(parts):
        col_idx = i % 2; row = i // 2
        x = x0 + col_idx * (cw + Inches(0.3))
        y = y0 + row * (ch + Inches(0.2))
        add_rounded(s, x, y, cw, ch, WHITE)
        add_rect(s, x, y, Inches(0.08), ch, col)
        # sign badge
        add_rounded(s, x + Inches(0.2), y + Inches(0.2),
                    Inches(0.5), Inches(0.5), col)
        add_text(s, x + Inches(0.2), y + Inches(0.2),
                 Inches(0.5), Inches(0.5), sign,
                 size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), y + Inches(0.1),
                 cw - Inches(1.0), Inches(0.35),
                 title, size=13, bold=True, color=NAVY)
        add_text(s, x + Inches(0.85), y + Inches(0.43),
                 cw - Inches(1.0), Inches(0.35),
                 formula, size=11, color=TEAL, font=FONT_MONO)
        add_text(s, x + Inches(0.85), y + Inches(0.72),
                 cw - Inches(1.0), Inches(0.35),
                 desc, size=10, color=GRAY_MED, italic=True)

    # Overall reward banner
    add_rounded(s, Inches(0.7), Inches(4.65), Inches(12.0), Inches(1.1), NAVY)
    add_text(s, Inches(1.0), Inches(4.78), Inches(11.5), Inches(0.3),
             "FINAL REWARD  =",
             size=11, bold=True, color=TEAL)
    add_text(s, Inches(1.0), Inches(5.05), Inches(11.5), Inches(0.65),
             "+ throughput_term  −  rtt_penalty  −  loss_penalty  +  stability_bonus",
             size=15, bold=True, color=WHITE, font=FONT_MONO)

    # Footer explanation
    add_text(s, Inches(0.7), Inches(6.0), Inches(12.0), Inches(0.9),
             "The reward is computed every 10 ms and clipped to [−2, 1.1]. "
             "Over training the agent learns that keeping RTT low near the BDP "
             "gives a better return than greedily maximising raw throughput.",
             size=12, italic=True, color=GRAY_MED)


# ── Slide 12 — Why SAC ───────────────────────────────────────────────────────
def build_why_sac():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Why Soft Actor-Critic (SAC)?",
                 kicker="The right RL algorithm for continuous, safety-sensitive control",
                 page_num=12, total=TOTAL_SLIDES)

    rows = [
        ["Action Space",       "Discrete only",   "Both",            "Continuous"],
        ["Sample Efficiency",  "Medium",          "Low (on-policy)", "High (off-policy)"],
        ["Exploration",        "ε-greedy",        "Stochastic",      "Entropy bonus"],
        ["Training Stability", "Sensitive",       "Stable",          "Very stable"],
    ]
    # Color the SAC column green for winners
    col_colors = {3: {0: GREEN, 1: GREEN, 2: GREEN, 3: GREEN}}
    add_table(s, Inches(0.7), Inches(1.65), Inches(12.0), Inches(2.8),
              ["Aspect", "DQN", "PPO", "SAC ✓"], rows,
              first_col_bold=True, col_colors=col_colors,
              header_size=12, body_size=12)

    # Mark SAC header cell green
    # (the above tables cannot override header per-column easily without XML)

    # Why SAC fits — 3 tiles
    y_tiles = Inches(4.9)
    tiles = [
        ("Continuous Actions",
         "cwnd factor ∈ [0.8, 1.2] — DQN cannot represent this."),
        ("Off-Policy Replay",
         "Every experience is re-used many times — precious in simulation."),
        ("Entropy Exploration",
         "Keeps the policy curious; prevents premature convergence."),
    ]
    tx = Inches(0.7); tw = Inches(3.95); th = Inches(1.6); gap = Inches(0.15)
    for i, (t, d) in enumerate(tiles):
        x = tx + i * (tw + gap)
        add_rounded(s, x, y_tiles, tw, th, WHITE)
        add_rect(s, x, y_tiles, Inches(0.08), th, TEAL)
        add_text(s, x + Inches(0.2), y_tiles + Inches(0.2),
                 tw - Inches(0.4), Inches(0.4),
                 t, size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.2), y_tiles + Inches(0.65),
                 tw - Inches(0.4), Inches(0.9),
                 d, size=11, color=GRAY_DARK)


# ── Slide 13 — SAC Algorithm ─────────────────────────────────────────────────
def build_sac_overview():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "SAC Algorithm — Actor & Critic",
                 kicker="Two neural networks cooperate to learn a good policy",
                 page_num=13, total=TOTAL_SLIDES)

    # Actor card
    ax = Inches(0.7); ay = Inches(1.65); aw = Inches(5.9); ah = Inches(3.5)
    add_rounded(s, ax, ay, aw, ah, WHITE)
    add_rect(s, ax, ay, aw, Inches(0.55), TEAL)
    add_text(s, ax + Inches(0.3), ay, aw, Inches(0.55),
             "ACTOR  (π)   —   “What to do”",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, ax + Inches(0.3), ay + Inches(0.8),
                aw - Inches(0.6), ah - Inches(1.0), [
        ("Input: ",  "current state (6 features)"),
        ("Output: ", "distribution over actions (Gaussian)"),
        ("Trained: ", "to produce high-Q, high-entropy actions"),
        ("At inference: ", "takes the mean (deterministic)"),
    ], size=12)

    # Critic card
    cx = Inches(6.8); cy = ay
    add_rounded(s, cx, cy, aw, ah, WHITE)
    add_rect(s, cx, cy, aw, Inches(0.55), NAVY)
    add_text(s, cx + Inches(0.3), cy, aw, Inches(0.55),
             "CRITIC  (Q)   —   “How good”",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, cx + Inches(0.3), cy + Inches(0.8),
                aw - Inches(0.6), ah - Inches(1.0), [
        ("Input: ",  "state + action"),
        ("Output: ", "Q-value  (expected long-term return)"),
        ("Trained: ", "to minimise Bellman-equation error"),
        ("Twin Qs: ", "two critics reduce over-estimation bias"),
    ], size=12)

    # Objective banner
    add_rounded(s, Inches(0.7), Inches(5.35), Inches(12.0), Inches(1.45), NAVY)
    add_text(s, Inches(1.0), Inches(5.48), Inches(11.5), Inches(0.3),
             "SAC OBJECTIVE  —  Maximum-Entropy RL",
             size=11, bold=True, color=ORANGE)
    add_text(s, Inches(1.0), Inches(5.78), Inches(11.5), Inches(0.55),
             "J(π) = E [ Σₜ  rₜ  +  α · H(π(· | sₜ)) ]",
             size=18, bold=True, color=WHITE, font=FONT_MONO)
    add_text(s, Inches(1.0), Inches(6.35), Inches(11.5), Inches(0.4),
             "Reward  +  entropy bonus  →  stable learning and principled exploration",
             size=12, color=RGBColor(0xB8, 0xD4, 0xF5), italic=True)


# ── Slide 14 — Neural Network Architecture ───────────────────────────────────
def build_nn_arch():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Neural Network Architecture",
                 kicker="Same backbone used for both actor and critic",
                 page_num=14, total=TOTAL_SLIDES)

    # Vertical stack of layer cards
    layers = [
        ("Input Layer",     "6 normalized features",       BLUE,  "6"),
        ("Dense  +  ReLU",  "Fully connected, 256 units",  TEAL,  "256"),
        ("Dense  +  ReLU",  "Fully connected, 256 units",  TEAL,  "256"),
        ("Dense  +  ReLU",  "Fully connected, 128 units",  TEAL,  "128"),
        ("Output Layer",    "Action mean  +  log std-dev", ORANGE, "2"),
    ]
    lx = Inches(1.5); ly = Inches(1.65); lw = Inches(5.5); lh = Inches(0.8)
    gap = Inches(0.25)
    for i, (name, desc, col, sz) in enumerate(layers):
        y = ly + i * (lh + gap)
        add_rounded(s, lx, y, lw, lh, WHITE)
        add_rect(s, lx, y, Inches(0.12), lh, col)
        add_text(s, lx + Inches(0.3), y + Inches(0.1),
                 Inches(2.7), Inches(0.35),
                 name, size=13, bold=True, color=NAVY)
        add_text(s, lx + Inches(0.3), y + Inches(0.42),
                 Inches(2.7), Inches(0.35),
                 desc, size=10, color=GRAY_MED)
        # size badge
        add_rounded(s, lx + lw - Inches(1.1), y + Inches(0.2),
                    Inches(0.9), Inches(0.4), col)
        add_text(s, lx + lw - Inches(1.1), y + Inches(0.2),
                 Inches(0.9), Inches(0.4), sz,
                 size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # arrow between layers
        if i < len(layers) - 1:
            add_arrow(s, lx + lw / 2, y + lh,
                      lx + lw / 2, y + lh + gap,
                      color=GRAY_MED, weight=1.5)

    # Right column — spec
    rx = Inches(8.3); ry = Inches(1.65)
    add_text(s, rx, ry, Inches(4.5), Inches(0.4),
             "MODEL SPECIFICATIONS", size=12, bold=True, color=ORANGE)
    specs = [
        ("Trainable parameters",   "≈ 200,000"),
        ("Activation",             "ReLU (hidden), tanh (output)"),
        ("Framework",              "PyTorch via stable-baselines3"),
        ("Optimizer",              "Adam,  lr = 3 × 10⁻⁴"),
        ("Discount factor γ",      "0.99"),
        ("Replay buffer size",     "1,000,000"),
        ("Mini-batch size",        "256"),
    ]
    for i, (k, v) in enumerate(specs):
        y = ry + Inches(0.5 + i * 0.55)
        add_rounded(s, rx, y, Inches(4.5), Inches(0.5), WHITE)
        add_text(s, rx + Inches(0.2), y, Inches(2.3), Inches(0.5),
                 k, size=11, color=GRAY_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + Inches(2.5), y, Inches(1.9), Inches(0.5),
                 v, size=11, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


# ── Slide 15 — System Architecture ───────────────────────────────────────────
def build_sys_arch():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "System Architecture",
                 kicker="Four cooperating layers, each with a specific job",
                 page_num=15, total=TOTAL_SLIDES)

    layers = [
        ("REACT DASHBOARD", "Live charts — RTT, Throughput, cwnd, Reward",
         "Port 3000  •  Recharts  •  WebSocket client", RGBColor(0x61, 0xDB, 0xFB)),
        ("SPRING BOOT BACKEND", "REST API  •  H2 persistence  •  STOMP broadcast",
         "Port 8080  •  Java 17  •  /topic/metrics", ORANGE),
        ("PYTHON AGENT", "Loads SAC model  •  inference loop  •  posts metrics",
         "stable-baselines3  •  PyTorch  •  10 ms step", RGBColor(0x3B, 0x82, 0xF6)),
        ("NS-3 SIMULATOR", "Dumbbell topology  •  TCP state machine  •  tracing",
         "C++ 17  •  ns-3.35  •  ns3-ai shared memory", RGBColor(0x6B, 0x5B, 0x95)),
    ]
    protocols = [
        "STOMP / WebSocket  (real-time push)",
        "HTTP POST  /api/metrics  (JSON)",
        "Shared memory  (ns3-ai, SHM id 2333)",
    ]

    lx = Inches(1.0); lw = Inches(11.3); lh = Inches(1.0)
    ly_start = Inches(1.65)
    gap = Inches(0.25)

    for i, (title, desc, stack, col) in enumerate(layers):
        y = ly_start + i * (lh + gap)
        add_rounded(s, lx, y, lw, lh, WHITE)
        add_rect(s, lx, y, Inches(0.15), lh, col)
        add_text(s, lx + Inches(0.35), y + Inches(0.1),
                 Inches(4.5), Inches(0.35),
                 title, size=14, bold=True, color=NAVY)
        add_text(s, lx + Inches(0.35), y + Inches(0.48),
                 Inches(8), Inches(0.3),
                 desc, size=11, color=GRAY_DARK)
        add_text(s, lx + Inches(0.35), y + Inches(0.75),
                 Inches(8), Inches(0.25),
                 stack, size=9, color=GRAY_MED, italic=True)

        # protocol label between layers
        if i < len(layers) - 1:
            proto_y = y + lh + Inches(0.02)
            add_text(s, lx + Inches(4.5), proto_y, Inches(6.5),
                     Inches(0.25), "▲  " + protocols[i],
                     size=10, bold=True, color=TEAL, italic=True)


# ── Slide 16 — Network Topology ──────────────────────────────────────────────
def build_topology():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Network Topology",
                 kicker="Classical dumbbell — both algorithms share the bottleneck",
                 page_num=16, total=TOTAL_SLIDES)

    # Canvas
    cx, cy, cw, ch = Inches(0.7), Inches(1.65), Inches(12.0), Inches(3.8)
    add_rounded(s, cx, cy, cw, ch, WHITE)

    # Nodes positions
    sac_s_x = cx + Inches(0.5); sac_s_y = cy + Inches(0.8)
    cub_s_x = cx + Inches(0.5); cub_s_y = cy + Inches(2.4)
    r0_x    = cx + Inches(4.0); r0_y    = cy + Inches(1.6)
    r1_x    = cx + Inches(7.7); r1_y    = cy + Inches(1.6)
    sac_r_x = cx + Inches(10.5); sac_r_y = cy + Inches(0.8)
    cub_r_x = cx + Inches(10.5); cub_r_y = cy + Inches(2.4)

    nw, nh = Inches(2.0), Inches(0.9)

    # Senders
    add_rounded(s, sac_s_x, sac_s_y, nw, nh, TEAL)
    add_text(s, sac_s_x, sac_s_y, nw, nh,
             "SAC Sender\n(Node 2)", size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rounded(s, cub_s_x, cub_s_y, nw, nh, BLUE)
    add_text(s, cub_s_x, cub_s_y, nw, nh,
             "CUBIC Sender\n(Node 3)", size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Routers
    add_rounded(s, r0_x, r0_y, nw, nh, NAVY)
    add_text(s, r0_x, r0_y, nw, nh, "Router 0",
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rounded(s, r1_x, r1_y, nw, nh, NAVY)
    add_text(s, r1_x, r1_y, nw, nh, "Router 1",
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Receivers
    add_rounded(s, sac_r_x, sac_r_y, nw, nh, GREEN)
    add_text(s, sac_r_x, sac_r_y, nw, nh,
             "SAC Sink\n(Node 4)", size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rounded(s, cub_r_x, cub_r_y, nw, nh, GREEN)
    add_text(s, cub_r_x, cub_r_y, nw, nh,
             "CUBIC Sink\n(Node 5)", size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Links
    # senders → R0
    add_arrow(s, sac_s_x + nw, sac_s_y + nh/2,
              r0_x, r0_y + nh/2, color=GRAY_DARK, weight=2.0)
    add_arrow(s, cub_s_x + nw, cub_s_y + nh/2,
              r0_x, r0_y + nh/2, color=GRAY_DARK, weight=2.0)
    # bottleneck R0 → R1 (bold red)
    add_arrow(s, r0_x + nw, r0_y + nh/2,
              r1_x, r1_y + nh/2, color=RED, weight=4.0)
    # bottleneck label
    add_text(s, r0_x + nw + Inches(0.15), r0_y - Inches(0.1),
             r1_x - (r0_x + nw) - Inches(0.3), Inches(0.4),
             "2 Mbps  /  20 ms  /  FqCoDel 100p",
             size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)
    # R1 → receivers
    add_arrow(s, r1_x + nw, r1_y + nh/2,
              sac_r_x, sac_r_y + nh/2, color=GRAY_DARK, weight=2.0)
    add_arrow(s, r1_x + nw, r1_y + nh/2,
              cub_r_x, cub_r_y + nh/2, color=GRAY_DARK, weight=2.0)
    # access link labels
    add_text(s, sac_s_x + nw + Inches(0.1), sac_s_y + nh/2 - Inches(0.3),
             Inches(1.8), Inches(0.3),
             "10 Mbps / 20 ms", size=9, color=GRAY_MED, italic=True)
    add_text(s, r1_x + nw + Inches(0.1), sac_r_y + nh/2 - Inches(0.3),
             Inches(1.8), Inches(0.3),
             "10 Mbps / 20 ms", size=9, color=GRAY_MED, italic=True)

    # Params strip
    params = [
        ("Bottleneck", "2 Mbps / 20 ms"),
        ("Access link", "10 Mbps / 20 ms"),
        ("Base RTT",   "≈ 80 ms"),
        ("BDP",        "≈ 2,500 bytes"),
    ]
    sx = Inches(0.7); sy = Inches(5.75)
    pw = Inches(2.95); ph = Inches(1.15); g = Inches(0.15)
    for i, (k, v) in enumerate(params):
        x = sx + i * (pw + g)
        add_rounded(s, x, sy, pw, ph, WHITE)
        add_rect(s, x, sy, Inches(0.08), ph, ORANGE)
        add_text(s, x + Inches(0.2), sy + Inches(0.15),
                 pw - Inches(0.4), Inches(0.35),
                 k, size=11, bold=True, color=TEAL)
        add_text(s, x + Inches(0.2), sy + Inches(0.5),
                 pw - Inches(0.4), Inches(0.55),
                 v, size=16, bold=True, color=NAVY)


# ── Slide 17 — Shared Memory ─────────────────────────────────────────────────
def build_shm():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Shared Memory Communication",
                 kicker="Why ns-3 and Python must talk at sub-millisecond speed",
                 page_num=17, total=TOTAL_SLIDES)

    # Latency comparison table
    rows = [
        ["TCP Socket (localhost)", "~10 µs", "Marginal"],
        ["gRPC",                   "~50 µs", "Acceptable"],
        ["Shared Memory (ns3-ai)", "~100 ns", "Optimal  ✓"],
    ]
    col_colors = {2: {0: GRAY_DARK, 1: GRAY_DARK, 2: GREEN}}
    add_table(s, Inches(0.7), Inches(1.65), Inches(6.5), Inches(2.0),
              ["Method", "Round-trip Latency", "Fit for 10 ms steps"],
              rows, first_col_bold=True, col_colors=col_colors,
              header_size=12, body_size=12)

    # Protocol flow (right)
    px = Inches(7.5); py = Inches(1.65)
    add_text(s, px, py, Inches(5.3), Inches(0.4),
             "LOCK-STEP EXCHANGE (every 10 ms)",
             size=12, bold=True, color=ORANGE)
    steps = [
        ("1", "C++ writes observation", "struct sTcpRlInferenceEnv", TEAL),
        ("2", "Python reads obs, runs SAC", "model.predict(state)",  BLUE),
        ("3", "Python writes action", "struct TcpRlInferenceAct",    ORANGE),
        ("4", "C++ reads, updates cwnd", "tcb->m_cWnd = new_cwnd",    GREEN),
    ]
    for i, (n, t, code, col) in enumerate(steps):
        y = py + Inches(0.5 + i * 0.65)
        add_rounded(s, px, y, Inches(5.3), Inches(0.55), WHITE)
        add_rounded(s, px + Inches(0.1), y + Inches(0.08),
                    Inches(0.4), Inches(0.4), col)
        add_text(s, px + Inches(0.1), y + Inches(0.08),
                 Inches(0.4), Inches(0.4), n,
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, px + Inches(0.65), y + Inches(0.05),
                 Inches(3.2), Inches(0.3), t,
                 size=11, bold=True, color=NAVY)
        add_text(s, px + Inches(0.65), y + Inches(0.3),
                 Inches(4.4), Inches(0.3), code,
                 size=9, color=GRAY_MED, font=FONT_MONO)

    # Code block showing the struct
    add_code_block(s, Inches(0.7), Inches(4.2), Inches(6.5), Inches(2.5),
                   "struct sTcpRlInferenceEnv {\n"
                   "    uint32_t cWnd;\n"
                   "    int64_t  rtt_us;\n"
                   "    double   throughput;\n"
                   "    uint32_t packetLoss;\n"
                   "    uint32_t segmentSize;\n"
                   "    uint32_t bytesInFlight;\n"
                   "    /* ... */\n"
                   "} Packed;",
                   title="Packed C struct shared with Python (ctypes)")

    # Note banner
    add_rounded(s, Inches(0.7), Inches(6.85), Inches(12.0), Inches(0.55), NAVY)
    add_text(s, Inches(1.0), Inches(6.85), Inches(11.5), Inches(0.55),
             "Synchronisation: POSIX condition variable → "
             "producer-consumer rendezvous every step",
             size=11, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)


# ── Slide 18 — Training Pipeline ─────────────────────────────────────────────
def build_training():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Training Pipeline",
                 kicker="Two-phase curriculum with reward refinement",
                 page_num=18, total=TOTAL_SLIDES)

    # Phase cards
    y0 = Inches(1.65); ph = Inches(2.0)
    # Phase 1
    add_rounded(s, Inches(0.7), y0, Inches(6.0), ph, WHITE)
    add_rect(s, Inches(0.7), y0, Inches(6.0), Inches(0.55), BLUE)
    add_text(s, Inches(0.9), y0, Inches(6), Inches(0.55),
             "Phase 1 — Initial Training",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(0.9), y0 + Inches(0.75),
                Inches(5.6), ph - Inches(0.85), [
        ("Steps: ",        "2,000,000"),
        ("RTT penalty: ",  "linear  (× 12)"),
        ("Goal: ",         "learn to maximise throughput"),
        ("Outcome: ",      "basic policy that keeps the link busy"),
    ], size=12)

    # Phase 2
    add_rounded(s, Inches(7.05), y0, Inches(6.0), ph, WHITE)
    add_rect(s, Inches(7.05), y0, Inches(6.0), Inches(0.55), ORANGE)
    add_text(s, Inches(7.25), y0, Inches(6), Inches(0.55),
             "Phase 2 — Fine-Tuning",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(7.25), y0 + Inches(0.75),
                Inches(5.6), ph - Inches(0.85), [
        ("Steps: ",        "200,000"),
        ("RTT penalty: ",  "quadratic  (× 40)"),
        ("Goal: ",         "reduce buffer bloat, stabilise cwnd"),
        ("Outcome: ",      "latency-aware, production-ready policy"),
    ], size=12)

    # Hyper-parameter strip
    add_text(s, Inches(0.7), Inches(3.95), Inches(12), Inches(0.4),
             "KEY TRAINING HYPER-PARAMETERS",
             size=12, bold=True, color=TEAL)
    hparams = [
        ("Learning rate", "3 × 10⁻⁴"),
        ("Replay buffer", "1,000,000"),
        ("Batch size",    "256"),
        ("Discount γ",    "0.99"),
        ("Step interval", "10 ms (sim)"),
    ]
    sx = Inches(0.7); sy = Inches(4.4)
    pw = Inches(2.4); ph2 = Inches(1.2); g = Inches(0.15)
    for i, (k, v) in enumerate(hparams):
        x = sx + i * (pw + g)
        add_rounded(s, x, sy, pw, ph2, WHITE)
        add_rect(s, x, sy, pw, Inches(0.08), TEAL)
        add_text(s, x + Inches(0.1), sy + Inches(0.2),
                 pw - Inches(0.2), Inches(0.4),
                 k, size=10, bold=True, color=GRAY_MED,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), sy + Inches(0.6),
                 pw - Inches(0.2), Inches(0.55),
                 v, size=16, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)

    # Insight banner
    add_rounded(s, Inches(0.7), Inches(5.9), Inches(12.0), Inches(1.1), NAVY)
    add_text(s, Inches(1.0), Inches(6.02), Inches(11.5), Inches(0.3),
             "WHY TWO PHASES?",
             size=11, bold=True, color=ORANGE)
    add_text(s, Inches(1.0), Inches(6.32), Inches(11.5), Inches(0.7),
             "Phase 1 teaches the agent to USE the link.  "
             "Phase 2 teaches it to USE IT GENTLY.",
             size=14, bold=True, color=WHITE, italic=True)


# ── Slide 19 — Experimental Setup ────────────────────────────────────────────
def build_exp_setup():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Experimental Setup",
                 kicker="Reproducible head-to-head evaluation",
                 page_num=19, total=TOTAL_SLIDES)

    # Left — methodology
    add_text(s, Inches(0.7), Inches(1.65), Inches(6), Inches(0.4),
             "METHODOLOGY", size=12, bold=True, color=ORANGE)
    methodology = [
        ("Simulation duration: ",  "200 seconds per run"),
        ("Random seed: ",          "42  (reproducibility)"),
        ("Runs per algorithm: ",   "3 × averaged"),
        ("Metrics captured: ",     "Throughput, RTT, Loss, Reward"),
        ("Sampling rate: ",        "every 10 ms (SAC), 100 ms (CUBIC)"),
    ]
    add_bullets(s, Inches(0.7), Inches(2.1), Inches(6), Inches(3.5),
                methodology, size=13)

    # Right — baselines
    add_text(s, Inches(7.0), Inches(1.65), Inches(5.7), Inches(0.4),
             "ALGORITHMS COMPARED", size=12, bold=True, color=ORANGE)
    baselines = [
        ("DRL-SAC",      "Our proposed learning-based controller", TEAL),
        ("TCP CUBIC",    "Linux default — cubic growth",           BLUE),
        ("TCP NewReno",  "Classical AIMD baseline",                 NAVY),
        ("TCP BBR",      "Google's model-based controller",         ORANGE),
    ]
    for i, (name, desc, col) in enumerate(baselines):
        y = Inches(2.1 + i * 0.8)
        add_rounded(s, Inches(7.0), y, Inches(5.7), Inches(0.7), WHITE)
        add_rect(s, Inches(7.0), y, Inches(0.1), Inches(0.7), col)
        add_text(s, Inches(7.2), y + Inches(0.08),
                 Inches(5.3), Inches(0.3),
                 name, size=13, bold=True, color=NAVY)
        add_text(s, Inches(7.2), y + Inches(0.4),
                 Inches(5.3), Inches(0.3),
                 desc, size=10, color=GRAY_MED, italic=True)

    # Footer banner
    add_rounded(s, Inches(0.7), Inches(5.75), Inches(12.0), Inches(1.25), NAVY)
    add_text(s, Inches(1.0), Inches(5.88), Inches(11.5), Inches(0.3),
             "EVALUATION SCENARIOS",
             size=11, bold=True, color=TEAL)
    add_text(s, Inches(1.0), Inches(6.18), Inches(11.5), Inches(0.4),
             "①  Baseline (no packet loss)       "
             "②  With 0.1 % random packet loss",
             size=13, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(6.55), Inches(11.5), Inches(0.4),
             "Same topology, same traffic pattern — only the congestion-control "
             "algorithm changes.",
             size=11, color=RGBColor(0xB8, 0xD4, 0xF5), italic=True)


# ── Slide 20 — Dashboard Screenshot placeholder 1 ────────────────────────────
def build_ss_placeholder(n, title, subtitle):
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, title, kicker=subtitle,
                 page_num=n, total=TOTAL_SLIDES)
    # Placeholder frame
    fx, fy, fw, fh = Inches(0.7), Inches(1.65), Inches(12.0), Inches(5.0)
    ph = slide_make_placeholder(s, fx, fy, fw, fh)
    # Caption
    add_text(s, Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.35),
             "⌘  Replace this placeholder with a screenshot of the running dashboard.",
             size=11, color=GRAY_MED, italic=True, align=PP_ALIGN.CENTER)


def slide_make_placeholder(s, x, y, w, h):
    # dashed border rounded rect
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.03
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = TEAL
    shp.line.width = Pt(2.0)
    # dashed style
    ln = shp.line._get_or_add_ln()
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")
    shp.shadow.inherit = False
    # central icon/text
    add_text(s, x, y, w, h,
             "🖼   PASTE DASHBOARD SCREENSHOT HERE",
             size=20, bold=True, color=TEAL,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, y + h/2 + Inches(0.4), w, Inches(0.4),
             "(right-click → Change Picture)",
             size=11, color=GRAY_MED, italic=True,
             align=PP_ALIGN.CENTER)
    return shp


# ── Slide 22 — Technology Stack ──────────────────────────────────────────────
def build_tech_stack():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Technology Stack",
                 kicker="One purpose-built tool per layer",
                 page_num=22, total=TOTAL_SLIDES)

    stack = [
        ("Simulation",      "ns-3 (C++)",               "Packet-level network simulation",          RGBColor(0x6B, 0x5B, 0x95)),
        ("IPC Bridge",      "ns3-ai (shared memory)",   "Sub-millisecond C++ ↔ Python",             RGBColor(0x8E, 0x44, 0xAD)),
        ("RL Agent",        "Python + stable-baselines3", "SAC model training & inference",         RGBColor(0x3B, 0x82, 0xF6)),
        ("Backend",         "Spring Boot + H2",         "REST API, persistence, broadcasting",      ORANGE),
        ("Frontend",        "React + Recharts",         "Live visualisation components",            RGBColor(0x61, 0xDB, 0xFB)),
        ("Realtime Protocol", "STOMP over WebSocket",   "Pub / sub streaming to the browser",       TEAL),
    ]
    # Build as an icon grid
    x0 = Inches(0.7); y0 = Inches(1.65)
    cw = Inches(6.1); ch = Inches(1.55); gap = Inches(0.15)
    for i, (role, name, desc, col) in enumerate(stack):
        col_idx = i % 2; row = i // 2
        x = x0 + col_idx * (cw + gap); y = y0 + row * (ch + gap)
        add_rounded(s, x, y, cw, ch, WHITE)
        add_rect(s, x, y, Inches(0.12), ch, col)
        # icon circle with first letter
        add_rounded(s, x + Inches(0.3), y + Inches(0.25),
                    Inches(1.05), Inches(1.05), col)
        add_text(s, x + Inches(0.3), y + Inches(0.25),
                 Inches(1.05), Inches(1.05),
                 name[0], size=28, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.55), y + Inches(0.25),
                 cw - Inches(1.7), Inches(0.35),
                 role, size=10, bold=True, color=TEAL)
        add_text(s, x + Inches(1.55), y + Inches(0.55),
                 cw - Inches(1.7), Inches(0.4),
                 name, size=15, bold=True, color=NAVY)
        add_text(s, x + Inches(1.55), y + Inches(0.97),
                 cw - Inches(1.7), Inches(0.5),
                 desc, size=11, color=GRAY_MED, italic=True)


# ── Slide 23 — Key Contributions ─────────────────────────────────────────────
def build_contributions():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Key Contributions",
                 kicker="What this project uniquely delivers",
                 page_num=23, total=TOTAL_SLIDES)

    items = [
        ("End-to-End Integration",
         "Working SAC ↔ ns-3 pipeline via ns3-ai shared memory at 10 ms steps."),
        ("Two-Phase Training Pipeline",
         "Initial throughput-focused training followed by latency-focused fine-tuning."),
        ("Loss-Robust Policy",
         "Learns to distinguish congestion loss from random loss — a capability that classical algorithms lack."),
        ("Full-Stack Monitoring",
         "Spring Boot backend + React dashboard stream live metrics for both SAC and CUBIC side-by-side."),
        ("Reproducible Research",
         "Fixed seeds, documented hyper-parameters, and a consistent comparison framework."),
    ]
    y = Inches(1.65)
    for i, (title, desc) in enumerate(items):
        h = Inches(0.95)
        add_rounded(s, Inches(0.7), y, Inches(12.0), h, WHITE)
        # number badge
        add_rounded(s, Inches(0.85), y + Inches(0.17),
                    Inches(0.6), Inches(0.6), TEAL)
        add_text(s, Inches(0.85), y + Inches(0.17),
                 Inches(0.6), Inches(0.6), f"{i+1:02d}",
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.6), y + Inches(0.13),
                 Inches(11.0), Inches(0.35),
                 title, size=14, bold=True, color=NAVY)
        add_text(s, Inches(1.6), y + Inches(0.48),
                 Inches(11.0), Inches(0.45),
                 desc, size=11, color=GRAY_DARK)
        y = y + h + Inches(0.1)


# ── Slide 24 — Limitations ───────────────────────────────────────────────────
def build_limitations():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Limitations",
                 kicker="Honest account of what this system does not yet solve",
                 page_num=24, total=TOTAL_SLIDES)

    items = [
        ("Single Topology",        "Trained only on a 2 Mbps / 80 ms dumbbell — "
                                   "generalisation to other networks is untested."),
        ("Multi-Flow Fairness",    "Behaviour against many competing flows has not been "
                                   "systematically evaluated."),
        ("Residual Bufferbloat",   "RTT is still slightly above the physical baseline "
                                   "under heavy load."),
        ("Simulation-Only",        "The agent has not been validated on a real Linux "
                                   "kernel or hardware NIC."),
        ("Hand-Tuned Reward",      "Reward coefficients are fixed constants chosen "
                                   "by the authors, not learned."),
        ("No Online Adaptation",   "Inference uses a frozen model — it cannot "
                                   "self-correct after deployment-time drift."),
    ]
    x0 = Inches(0.7); y0 = Inches(1.65)
    cw = Inches(6.0); ch = Inches(1.5); gap = Inches(0.2)
    for i, (t, d) in enumerate(items):
        col = i % 2; row = i // 2
        x = x0 + col * (cw + gap)
        y = y0 + row * (ch + gap)
        add_rounded(s, x, y, cw, ch, WHITE)
        add_rect(s, x, y, Inches(0.1), ch, RED)
        add_text(s, x + Inches(0.25), y + Inches(0.2),
                 cw - Inches(0.4), Inches(0.4),
                 "⚠  " + t, size=13, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), y + Inches(0.65),
                 cw - Inches(0.4), Inches(0.8),
                 d, size=11, color=GRAY_DARK)


# ── Slide 25 — Future Work ───────────────────────────────────────────────────
def build_future():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Future Work",
                 kicker="Short-term improvements and long-term research directions",
                 page_num=25, total=TOTAL_SLIDES)

    immediate = [
        ("Diverse topologies",  "Train across varying bandwidth, delay, queue sizes"),
        ("Multi-flow evaluation", "Quantify fairness vs BBRv2, CUBIC, Reno mixes"),
        ("Real traces",         "Replay captured WAN / 4G packet traces"),
    ]
    advanced = [
        ("Recurrent / Transformer policy", "Capture long-range temporal patterns"),
        ("Curriculum learning",            "Easy networks first, then harder ones"),
        ("Safe RL constraints",            "Bound worst-case RTT during exploration"),
        ("Linux kernel deployment",        "Export via eBPF / TCP_CONG hook"),
    ]

    # left card — immediate
    cx, cy = Inches(0.7), Inches(1.65); cw, ch = Inches(6.0), Inches(5.2)
    add_rounded(s, cx, cy, cw, ch, WHITE)
    add_rect(s, cx, cy, cw, Inches(0.55), TEAL)
    add_text(s, cx + Inches(0.3), cy, cw, Inches(0.55),
             "Immediate Improvements",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    for i, (t, d) in enumerate(immediate):
        y = cy + Inches(0.85 + i * 1.35)
        add_rounded(s, cx + Inches(0.3), y, cw - Inches(0.6), Inches(1.15),
                    BG)
        add_text(s, cx + Inches(0.5), y + Inches(0.15),
                 cw - Inches(1.0), Inches(0.35),
                 "▸ " + t, size=13, bold=True, color=NAVY)
        add_text(s, cx + Inches(0.5), y + Inches(0.55),
                 cw - Inches(1.0), Inches(0.55),
                 d, size=11, color=GRAY_MED, italic=True)

    # right card — advanced
    cx2 = Inches(7.05)
    add_rounded(s, cx2, cy, cw, ch, WHITE)
    add_rect(s, cx2, cy, cw, Inches(0.55), ORANGE)
    add_text(s, cx2 + Inches(0.3), cy, cw, Inches(0.55),
             "Advanced Research Extensions",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    for i, (t, d) in enumerate(advanced):
        y = cy + Inches(0.85 + i * 1.0)
        add_rounded(s, cx2 + Inches(0.3), y, cw - Inches(0.6), Inches(0.85),
                    BG)
        add_text(s, cx2 + Inches(0.5), y + Inches(0.1),
                 cw - Inches(1.0), Inches(0.35),
                 "▸ " + t, size=12, bold=True, color=NAVY)
        add_text(s, cx2 + Inches(0.5), y + Inches(0.42),
                 cw - Inches(1.0), Inches(0.4),
                 d, size=10, color=GRAY_MED, italic=True)


# ── Slide 26 — Conclusion ────────────────────────────────────────────────────
def build_conclusion():
    s = prs.slides.add_slide(BLANK); slide_background(s)
    slide_header(s, "Conclusion",
                 kicker="What we demonstrated and why it matters",
                 page_num=26, total=TOTAL_SLIDES)

    points = [
        "Deep Reinforcement Learning is a viable approach for TCP congestion control.",
        "A Soft Actor-Critic policy can match or exceed classical algorithms in normal conditions.",
        "The learned policy is significantly more robust when random packet loss is present.",
        "The agent learned to distinguish random loss from congestion — something classical algorithms cannot do.",
    ]
    y = Inches(1.65)
    for i, p in enumerate(points):
        h = Inches(0.85)
        add_rounded(s, Inches(0.7), y, Inches(12), h, WHITE)
        add_rect(s, Inches(0.7), y, Inches(0.1), h, GREEN)
        # check mark
        add_rounded(s, Inches(0.95), y + Inches(0.2),
                    Inches(0.45), Inches(0.45), GREEN)
        add_text(s, Inches(0.95), y + Inches(0.2),
                 Inches(0.45), Inches(0.45), "✓",
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.5), y + Inches(0.22),
                 Inches(11.0), Inches(0.45),
                 p, size=13, color=GRAY_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        y = y + h + Inches(0.1)

    # bottom banner
    add_rounded(s, Inches(0.7), Inches(5.9), Inches(12.0), Inches(1.2), NAVY)
    add_text(s, Inches(1.0), Inches(6.0), Inches(11.5), Inches(0.3),
             "BOTTOM LINE",
             size=11, bold=True, color=ORANGE)
    add_text(s, Inches(1.0), Inches(6.3), Inches(11.5), Inches(0.8),
             "Learning-based congestion control can adapt to network diversity "
             "that classical algorithms were never designed for.",
             size=15, bold=True, color=WHITE, italic=True)


# ── Slide 27 — Thank you / Q&A ───────────────────────────────────────────────
def build_thank_you():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, 0, SW, Inches(3.5), BLUE)
    add_rect(s, 0, Inches(3.5), SW, Inches(0.08), TEAL)

    # Big Thank You
    add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
             "Thank You", size=60, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.6),
             "Questions, critique, and discussion are welcome.",
             size=18, color=RGBColor(0xB8, 0xD4, 0xF5), italic=True)

    # Project title strip
    add_text(s, Inches(0.7), Inches(2.75), Inches(12), Inches(0.6),
             "Adaptive TCP Congestion Control using Deep Reinforcement Learning",
             size=16, bold=True, color=TEAL)

    # Team card
    tx = Inches(0.7); ty = Inches(4.0); tw = Inches(5.8); th = Inches(2.7)
    add_rounded(s, tx, ty, tw, th, WHITE)
    add_rect(s, tx, ty, Inches(0.08), th, ORANGE)
    add_text(s, tx + Inches(0.3), ty + Inches(0.2),
             tw - Inches(0.5), Inches(0.4),
             "TEAM", size=11, bold=True, color=TEAL)
    members = ["Jewel Bhannvadiya  (CE-065)",
               "Archi Patel  (CE-087)",
               "Henil Suchak  (CE-148)"]
    for i, m in enumerate(members):
        add_text(s, tx + Inches(0.3), ty + Inches(0.6 + i * 0.45),
                 tw - Inches(0.5), Inches(0.4),
                 "● " + m, size=14, bold=True, color=NAVY)
    add_text(s, tx + Inches(0.3), ty + Inches(2.1),
             tw - Inches(0.5), Inches(0.4),
             "Guided by Prof. Bhavika M. Gambhava",
             size=11, color=GRAY_MED, italic=True)

    # Q&A card
    qx = Inches(6.85); qy = Inches(4.0); qw = Inches(5.8); qh = Inches(2.7)
    add_rounded(s, qx, qy, qw, qh, TEAL)
    add_text(s, qx, qy, qw, qh, "Q & A",
             size=72, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ─────────────────────────────────────────────────────────────────────────────
# BUILD ALL SLIDES IN ORDER
# ─────────────────────────────────────────────────────────────────────────────

build_title()                  # 1
build_agenda()                 # 2
build_problem()                # 3
build_why_matters()            # 4
build_what_is_tcp()            # 5
build_congestion()             # 6
build_classical()              # 7
build_sawtooth()               # 8
build_why_rl()                 # 9
build_mdp()                    # 10
build_reward()                 # 11
build_why_sac()                # 12
build_sac_overview()           # 13
build_nn_arch()                # 14
build_sys_arch()               # 15
build_topology()               # 16
build_shm()                    # 17
build_training()               # 18
build_exp_setup()              # 19
build_ss_placeholder(20,
    "Live Dashboard — Overview",
    "Real-time monitoring of SAC and CUBIC side by side")    # 20
build_ss_placeholder(21,
    "Live Dashboard — SAC vs CUBIC Comparison",
    "Throughput, RTT and cwnd trajectories")                  # 21
build_tech_stack()             # 22
build_contributions()          # 23
build_limitations()            # 24
build_future()                 # 25
build_conclusion()             # 26
build_thank_you()              # 27


out_path = "Adaptive_TCP_Congestion_Control.pptx"
prs.save(out_path)
print(f"Wrote {out_path}  ({len(prs.slides)} slides)")
